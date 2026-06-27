import os
import re
import json
import logging
from typing import List, Dict, Any, Tuple, Optional

# Supported format parsers
from pypdf import PdfReader
# pyrefly: ignore [missing-import]
import docx
import pandas as pd
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from bs4 import BeautifulSoup
import markdown

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Universal document processor that extracts text from various formats,
    cleans it, and produces intelligent semantic chunks.
    """
    _MIN_CHUNK_CHARS = 800
    _MAX_CHUNK_CHARS = 2500
    _TARGET_WORDS = 300
    _OVERLAP_SENTENCES = 2

    _SECTION_HEADER_RE = re.compile(
        r"^(?:[A-Z][A-Z &/,\-]{2,}|.{3,60}:)\s*$", re.MULTILINE
    )

    @classmethod
    def process_document(cls, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a document based on its extension, and chunks it.
        Returns a list of chunk dictionaries with source metadata.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(filename)[1].lower()
        
        try:
            if ext == '.pdf':
                merged_text, page_ranges = cls._extract_pdf(file_path)
                doc_type = 'PDF'
            elif ext == '.docx':
                merged_text, page_ranges = cls._extract_docx(file_path)
                doc_type = 'DOCX'
            elif ext == '.txt':
                merged_text, page_ranges = cls._extract_txt(file_path)
                doc_type = 'TXT'
            elif ext == '.csv':
                merged_text, page_ranges = cls._extract_csv(file_path)
                doc_type = 'CSV'
            elif ext in ['.xlsx', '.xls']:
                merged_text, page_ranges = cls._extract_excel(file_path)
                doc_type = 'XLSX'
            elif ext == '.pptx':
                merged_text, page_ranges = cls._extract_pptx(file_path)
                doc_type = 'PPTX'
            elif ext == '.json':
                merged_text, page_ranges = cls._extract_json(file_path)
                doc_type = 'JSON'
            elif ext in ['.html', '.htm']:
                merged_text, page_ranges = cls._extract_html(file_path)
                doc_type = 'HTML'
            elif ext == '.md':
                merged_text, page_ranges = cls._extract_md(file_path)
                doc_type = 'MD'
            else:
                logger.warning(f"Unsupported file format: {ext}")
                return []
                
        except Exception as e:
            logger.error(f"Error extracting {filename}: {e}")
            return []

        if not merged_text.strip():
            logger.warning(f"No text extracted from {filename}")
            return []

        # Target chunk size calculation
        target_size = cls._auto_chunk_size(merged_text)
        sections = cls._split_into_sections(merged_text)
        
        all_chunks: List[Dict[str, Any]] = []
        chunk_index = 0

        for section_text, section_start_char in sections:
            sentences = cls._split_into_sentences(section_text)
            if not sentences:
                continue
                
            packed = cls._pack_sentences(
                sentences, target_size, section_start_char, page_ranges
            )
            
            for chunk_text, page_num in packed:
                all_chunks.append({
                    "text": chunk_text,
                    "source_file": filename,
                    "document_type": doc_type,
                    "page_number": page_num,
                    "chunk_index": chunk_index,
                })
                chunk_index += 1

        logger.info(f"Generated {len(all_chunks)} chunks from {filename}.")
        return all_chunks

    # ------------------------------------------------------------------ #
    # Format Extractors                                                  #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _extract_pdf(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        merged_text = ""
        page_ranges = []
        reader = PdfReader(file_path)
        
        for index, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            clean_text = DocumentProcessor._clean_page_text(page_text)
            if clean_text:
                sep = "\n\n" if merged_text else ""
                start_char = len(merged_text) + len(sep)
                merged_text += sep + clean_text
                page_ranges.append((start_char, len(merged_text), index + 1))
                
        return merged_text, page_ranges

    @staticmethod
    def _extract_docx(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text, [(0, len(text), 1)]

    @staticmethod
    def _extract_txt(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return text, [(0, len(text), 1)]

    @staticmethod
    def _extract_csv(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
        return text, [(0, len(text), None)]

    @staticmethod
    def _extract_excel(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        dfs = pd.read_excel(file_path, sheet_name=None)
        texts = []
        for sheet_name, df in dfs.items():
            texts.append(f"Sheet: {sheet_name}\n" + df.to_string(index=False))
        text = "\n\n".join(texts)
        return text, [(0, len(text), None)]

    @staticmethod
    def _extract_pptx(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        prs = Presentation(file_path)
        merged_text = ""
        page_ranges = []
        
        for index, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            clean_text = DocumentProcessor._clean_page_text("\n".join(slide_text))
            if clean_text:
                sep = "\n\n" if merged_text else ""
                start_char = len(merged_text) + len(sep)
                merged_text += sep + clean_text
                page_ranges.append((start_char, len(merged_text), index + 1))
                
        return merged_text, page_ranges

    @staticmethod
    def _extract_json(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        text = json.dumps(data, indent=2)
        return text, [(0, len(text), 1)]

    @staticmethod
    def _extract_html(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        text = soup.get_text(separator='\n', strip=True)
        return text, [(0, len(text), 1)]

    @staticmethod
    def _extract_md(file_path: str) -> Tuple[str, List[Tuple[int, int, int]]]:
        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text(separator='\n', strip=True)
        return text, [(0, len(text), 1)]

    # ------------------------------------------------------------------ #
    # Chunking Logic (Preserved)                                         #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _clean_page_text(raw: str) -> str:
        text = re.sub(r"-\n(\S)", r"\1", raw)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        lines = [line.strip() for line in text.split("\n")]
        return "\n".join(lines).strip()

    @classmethod
    def _auto_chunk_size(cls, text: str) -> int:
        words = len(text.split())
        if words < 500: return 1000
        elif words < 3000: return 1500
        elif words < 10000: return 2000
        return 2500

    @classmethod
    def _split_into_sections(cls, text: str) -> List[Tuple[str, int]]:
        headers = list(cls._SECTION_HEADER_RE.finditer(text))
        if not headers:
            return [(text, 0)]
        sections = []
        if headers[0].start() > 0:
            pre = text[: headers[0].start()].strip()
            if pre: sections.append((pre, 0))
        for i, match in enumerate(headers):
            start = match.start()
            end = headers[i + 1].start() if i + 1 < len(headers) else len(text)
            section_text = text[start:end].strip()
            if section_text:
                sections.append((section_text, start))
        return sections

    @staticmethod
    def _split_into_sentences(text: str) -> List[str]:
        text = re.sub(r"[\u2022\uf0b7]", ". ", text)
        sentences = re.split(r"(?<=[.!?])\s+", text)
        return [s.strip() for s in sentences if len(s.strip()) > 10]

    @staticmethod
    def get_page_number(char_pos: int, page_ranges: List[Tuple[int, int, int]]) -> Optional[int]:
        if not page_ranges:
            return None
        for start, end, page_num in page_ranges:
            if start <= char_pos <= end:
                return page_num
        for start, end, page_num in page_ranges:
            if char_pos <= end:
                return page_num
        return page_ranges[-1][2]

    @classmethod
    def _pack_sentences(
        cls, sentences: List[str], target_size: int, section_start_char: int,
        page_ranges: List[Tuple[int, int, int]]
    ) -> List[Tuple[str, int]]:
        results = []
        current = []
        current_len = 0
        
        for sentence in sentences:
            sentence_len = len(sentence)
            if current and current_len + sentence_len > target_size:
                chunk_text = " ".join(current)
                page_num = cls.get_page_number(section_start_char, page_ranges)
                results.append((chunk_text, page_num))
                overlap = current[-2:] if len(current) >= 2 else current
                current = overlap.copy()
                current_len = len(" ".join(current))
            current.append(sentence)
            current_len += sentence_len
            
        if current:
            chunk_text = " ".join(current)
            page_num = cls.get_page_number(section_start_char, page_ranges)
            results.append((chunk_text, page_num))
            
        return results
